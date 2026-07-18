"""
Parsers for Text, HTML, XML, JSON, Email, PowerPoint, and ZIP files.
Each parser extracts tabular data into the internal Table model.
"""
from __future__ import annotations

import csv
import email
import io
import json
import logging
import os
import re
import tempfile
import time
import zipfile
from pathlib import Path
from typing import Optional
from xml.etree import ElementTree

from backend.models import Cell, Table, ProcessingResult, ProcessingStatus
from backend.utils import detect_data_type, get_file_extension, cleanup_directory

logger = logging.getLogger(__name__)


# ── Text/RTF Parser ────────────────────────────────────────────────────────

class TextParser:
    """Parses plain text, CSV, and RTF files."""

    def parse(self, file_path: str | Path) -> ProcessingResult:
        result = ProcessingResult()
        result.filename = Path(file_path).name
        result.file_category = "text"
        result.page_count = 1
        start_time = time.time()

        try:
            result.status = ProcessingStatus.OCR_IN_PROGRESS
            text = self._read_file(str(file_path))

            if not text.strip():
                result.mark_failed("File is empty")
                return result

            # Try CSV parsing first
            table = self._try_csv(text)
            if not table:
                # Fall back to whitespace/tab delimited
                table = self._parse_structured_text(text)

            if table:
                table.source_engine = "text_parser"
                result.tables.append(table)

            result.ocr_engine_used = "text_parser"
            result.status = ProcessingStatus.READY_FOR_REVIEW
            result.processing_time = time.time() - start_time

        except Exception as e:
            result.mark_failed(f"Text parsing failed: {str(e)}")

        return result

    def _read_file(self, file_path: str) -> str:
        """Read file with encoding detection."""
        encodings = ["utf-8", "latin-1", "cp1252", "iso-8859-1", "ascii"]
        for enc in encodings:
            try:
                with open(file_path, "r", encoding=enc) as f:
                    return f.read()
            except (UnicodeDecodeError, UnicodeError):
                continue
        # Binary fallback
        with open(file_path, "rb") as f:
            return f.read().decode("utf-8", errors="replace")

    def _try_csv(self, text: str) -> Optional[Table]:
        """Try to parse as CSV."""
        try:
            dialect = csv.Sniffer().sniff(text[:4096])
            reader = csv.reader(io.StringIO(text), dialect)
            rows = list(reader)
            if len(rows) < 2:
                return None
            return self._rows_to_table(rows)
        except csv.Error:
            return None

    def _parse_structured_text(self, text: str) -> Optional[Table]:
        """Parse whitespace/tab delimited text."""
        lines = [line for line in text.split('\n') if line.strip()]
        if len(lines) < 2:
            return None

        rows: list[list[str]] = []
        for line in lines:
            parts = re.split(r'\t|\s{2,}|\|', line)
            parts = [p.strip() for p in parts if p.strip()]
            if parts:
                rows.append(parts)

        if len(rows) < 2:
            return None

        return self._rows_to_table(rows)

    def _rows_to_table(self, rows: list[list[str]]) -> Table:
        """Convert list of string rows to Table."""
        table = Table()
        table.has_borders = False

        for row_idx, row_data in enumerate(rows):
            cells: list[Cell] = []
            for col_idx, value in enumerate(row_data):
                cell = Cell(row=row_idx, col=col_idx)
                cell.raw_value = value
                cell.value = value
                cell.confidence = 1.0

                if value:
                    dt, parsed, fmt = detect_data_type(value)
                    cell.data_type = dt
                    if dt in ("number", "currency", "percentage"):
                        cell.value = parsed
                    cell.format_string = fmt

                if row_idx == 0:
                    cell.is_header = True
                    cell.font_bold = True

                cells.append(cell)
            table.cells.append(cells)

        if table.cells:
            table.headers = [str(c.value) for c in table.cells[0]]
            table.num_rows = len(table.cells)
            table.num_cols = max(len(r) for r in table.cells) if table.cells else 0

        return table


# ── HTML Parser ─────────────────────────────────────────────────────────────

class HTMLParser:
    """Parses HTML files and extracts tables."""

    def parse(self, file_path: str | Path) -> ProcessingResult:
        result = ProcessingResult()
        result.filename = Path(file_path).name
        result.file_category = "html"
        result.page_count = 1
        start_time = time.time()

        try:
            result.status = ProcessingStatus.OCR_IN_PROGRESS

            with open(str(file_path), "r", encoding="utf-8", errors="replace") as f:
                html_content = f.read()

            tables = self._extract_html_tables(html_content)
            result.tables = tables
            result.ocr_engine_used = "html_parser"
            result.status = ProcessingStatus.READY_FOR_REVIEW
            result.processing_time = time.time() - start_time

        except Exception as e:
            result.mark_failed(f"HTML parsing failed: {str(e)}")

        return result

    def _extract_html_tables(self, html: str) -> list[Table]:
        """Extract all <table> elements from HTML."""
        from html.parser import HTMLParser as StdHTMLParser

        tables: list[Table] = []
        current_table: list[list[str]] = []
        current_row: list[str] = []
        in_table = False
        in_cell = False
        cell_text = ""
        table_count = 0

        class TableExtractor(StdHTMLParser):
            nonlocal in_table, in_cell, cell_text, current_table, current_row, table_count

            def handle_starttag(self_, tag, attrs):
                nonlocal in_table, in_cell, cell_text, current_table, current_row
                tag = tag.lower()
                if tag == "table":
                    in_table = True
                    current_table = []
                elif tag == "tr" and in_table:
                    current_row = []
                elif tag in ("td", "th") and in_table:
                    in_cell = True
                    cell_text = ""

            def handle_endtag(self_, tag):
                nonlocal in_table, in_cell, cell_text, current_table, current_row, table_count
                tag = tag.lower()
                if tag in ("td", "th") and in_cell:
                    current_row.append(cell_text.strip())
                    in_cell = False
                elif tag == "tr" and in_table:
                    if current_row:
                        current_table.append(current_row)
                elif tag == "table" and in_table:
                    if current_table and len(current_table) > 1:
                        tbl = self._build_table_from_rows(current_table, table_count)
                        if tbl:
                            tables.append(tbl)
                            table_count += 1
                    in_table = False

            def handle_data(self_, data):
                nonlocal cell_text
                if in_cell:
                    cell_text += data

        parser = TableExtractor()
        parser.feed(html)
        return tables

    def _build_table_from_rows(self, rows: list[list[str]], index: int) -> Optional[Table]:
        """Build Table from parsed HTML rows."""
        table = Table()
        table.source_engine = "html_parser"
        table.title = f"Table {index + 1}"
        table.has_borders = True

        for row_idx, row_data in enumerate(rows):
            cells: list[Cell] = []
            for col_idx, value in enumerate(row_data):
                cell = Cell(row=row_idx, col=col_idx)
                cell.raw_value = value
                cell.value = value
                cell.confidence = 1.0

                if value:
                    dt, parsed, fmt = detect_data_type(value)
                    cell.data_type = dt
                    if dt in ("number", "currency", "percentage"):
                        cell.value = parsed
                    cell.format_string = fmt

                if row_idx == 0:
                    cell.is_header = True
                    cell.font_bold = True

                cells.append(cell)
            table.cells.append(cells)

        if table.cells:
            table.headers = [str(c.value) for c in table.cells[0]]
            table.num_rows = len(table.cells)
            table.num_cols = max(len(r) for r in table.cells) if table.cells else 0

        return table


# ── XML/JSON Parser ─────────────────────────────────────────────────────────

class XMLJSONParser:
    """Parses XML and JSON files containing tabular data."""

    def parse(self, file_path: str | Path) -> ProcessingResult:
        result = ProcessingResult()
        result.filename = Path(file_path).name
        result.file_category = "data"
        result.page_count = 1
        start_time = time.time()

        ext = get_file_extension(str(file_path))

        try:
            result.status = ProcessingStatus.OCR_IN_PROGRESS

            if ext == "json":
                tables = self._parse_json(str(file_path))
            elif ext == "xml":
                tables = self._parse_xml(str(file_path))
            else:
                result.mark_failed(f"Unsupported data format: {ext}")
                return result

            result.tables = tables
            result.ocr_engine_used = "data_parser"
            result.status = ProcessingStatus.READY_FOR_REVIEW
            result.processing_time = time.time() - start_time

        except Exception as e:
            result.mark_failed(f"Data parsing failed: {str(e)}")

        return result

    def _parse_json(self, file_path: str) -> list[Table]:
        """Parse JSON file. Handles arrays of objects and nested structures."""
        with open(file_path, "r", encoding="utf-8") as f:
            data = json.load(f)

        tables: list[Table] = []

        if isinstance(data, list) and data and isinstance(data[0], dict):
            # Array of objects → table
            table = self._dicts_to_table(data, "JSON Data")
            if table:
                tables.append(table)
        elif isinstance(data, dict):
            # Check for known GST/financial structures
            for key, value in data.items():
                if isinstance(value, list) and value and isinstance(value[0], dict):
                    table = self._dicts_to_table(value, key)
                    if table:
                        tables.append(table)

        return tables

    def _parse_xml(self, file_path: str) -> list[Table]:
        """Parse XML file into table format."""
        tree = ElementTree.parse(file_path)
        root = tree.getroot()

        tables: list[Table] = []

        # Find repeated elements (likely table rows)
        tag_counts: dict[str, int] = {}
        for child in root:
            tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            tag_counts[tag] = tag_counts.get(tag, 0) + 1

        for tag, count in tag_counts.items():
            if count < 2:
                continue

            elements = [e for e in root if (e.tag.split('}')[-1] if '}' in e.tag else e.tag) == tag]
            dicts = []
            for elem in elements:
                row_dict = {}
                for child in elem:
                    child_tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
                    row_dict[child_tag] = child.text or ""
                if row_dict:
                    dicts.append(row_dict)

            if dicts:
                table = self._dicts_to_table(dicts, tag)
                if table:
                    tables.append(table)

        return tables

    def _dicts_to_table(self, dicts: list[dict], title: str = "") -> Optional[Table]:
        """Convert list of dicts to Table."""
        if not dicts:
            return None

        # Collect all keys maintaining order
        keys: list[str] = []
        for d in dicts:
            for k in d:
                if k not in keys:
                    keys.append(k)

        table = Table()
        table.source_engine = "data_parser"
        table.title = title
        table.has_borders = True

        # Header row
        header_cells = []
        for col_idx, key in enumerate(keys):
            cell = Cell(row=0, col=col_idx)
            cell.value = key
            cell.raw_value = key
            cell.is_header = True
            cell.font_bold = True
            cell.confidence = 1.0
            header_cells.append(cell)
        table.cells.append(header_cells)
        table.headers = keys

        # Data rows
        for row_idx, d in enumerate(dicts):
            cells: list[Cell] = []
            for col_idx, key in enumerate(keys):
                cell = Cell(row=row_idx + 1, col=col_idx)
                value = str(d.get(key, ""))
                cell.raw_value = value
                cell.value = value
                cell.confidence = 1.0

                if value:
                    dt, parsed, fmt = detect_data_type(value)
                    cell.data_type = dt
                    if dt in ("number", "currency", "percentage"):
                        cell.value = parsed
                    cell.format_string = fmt

                cells.append(cell)
            table.cells.append(cells)

        table.num_rows = len(table.cells)
        table.num_cols = len(keys)
        return table


# ── Email Parser ────────────────────────────────────────────────────────────

class EmailParser:
    """Parses email files (.eml, .msg) extracting attachments and body tables."""

    def __init__(self, parent_processor=None):
        self.parent_processor = parent_processor

    def parse(self, file_path: str | Path) -> ProcessingResult:
        result = ProcessingResult()
        result.filename = Path(file_path).name
        result.file_category = "email"
        result.page_count = 1
        start_time = time.time()

        ext = get_file_extension(str(file_path))

        try:
            result.status = ProcessingStatus.OCR_IN_PROGRESS

            if ext == "eml":
                tables = self._parse_eml(str(file_path))
            elif ext == "msg":
                tables = self._parse_msg(str(file_path))
            else:
                result.mark_failed(f"Unsupported email format: {ext}")
                return result

            result.tables = tables
            result.ocr_engine_used = "email_parser"
            result.status = ProcessingStatus.READY_FOR_REVIEW
            result.processing_time = time.time() - start_time

        except Exception as e:
            result.mark_failed(f"Email parsing failed: {str(e)}")

        return result

    def _parse_eml(self, file_path: str) -> list[Table]:
        """Parse .eml email file."""
        with open(file_path, "rb") as f:
            msg = email.message_from_binary_file(f)

        tables: list[Table] = []

        # Extract HTML body tables
        for part in msg.walk():
            content_type = part.get_content_type()
            if content_type == "text/html":
                html = part.get_payload(decode=True)
                if html:
                    html_parser = HTMLParser()
                    html_tables = html_parser._extract_html_tables(
                        html.decode("utf-8", errors="replace")
                    )
                    tables.extend(html_tables)

        return tables

    def _parse_msg(self, file_path: str) -> list[Table]:
        """Parse .msg Outlook email file."""
        try:
            import extract_msg
            msg = extract_msg.Message(file_path)
            tables: list[Table] = []

            if msg.htmlBody:
                html_parser = HTMLParser()
                html_tables = html_parser._extract_html_tables(
                    msg.htmlBody.decode("utf-8", errors="replace")
                    if isinstance(msg.htmlBody, bytes) else msg.htmlBody
                )
                tables.extend(html_tables)

            msg.close()
            return tables
        except ImportError:
            logger.warning("extract-msg not installed, .msg parsing unavailable")
            return []


# ── PowerPoint Parser ───────────────────────────────────────────────────────

class PPTParser:
    """Parses PowerPoint files (.pptx) extracting tables from slides."""

    def parse(self, file_path: str | Path) -> ProcessingResult:
        result = ProcessingResult()
        result.filename = Path(file_path).name
        result.file_category = "powerpoint"
        start_time = time.time()

        try:
            result.status = ProcessingStatus.OCR_IN_PROGRESS

            from pptx import Presentation
            prs = Presentation(str(file_path))
            result.page_count = len(prs.slides)

            tables: list[Table] = []
            for slide_num, slide in enumerate(prs.slides, 1):
                for shape in slide.shapes:
                    if shape.has_table:
                        table = self._extract_pptx_table(shape.table, slide_num)
                        if table:
                            tables.append(table)

            result.tables = tables
            result.ocr_engine_used = "python-pptx"
            result.status = ProcessingStatus.READY_FOR_REVIEW
            result.processing_time = time.time() - start_time

        except Exception as e:
            result.mark_failed(f"PowerPoint parsing failed: {str(e)}")

        return result

    def _extract_pptx_table(self, pptx_table, slide_num: int) -> Optional[Table]:
        """Extract table from a PowerPoint table shape."""
        table = Table()
        table.source_engine = "python-pptx"
        table.page_number = slide_num
        table.title = f"Slide {slide_num} Table"

        for row_idx, row in enumerate(pptx_table.rows):
            cells: list[Cell] = []
            for col_idx, pptx_cell in enumerate(row.cells):
                cell = Cell(row=row_idx, col=col_idx)
                text = pptx_cell.text.strip()
                cell.raw_value = text
                cell.value = text
                cell.confidence = 1.0

                if text:
                    dt, parsed, fmt = detect_data_type(text)
                    cell.data_type = dt
                    if dt in ("number", "currency", "percentage"):
                        cell.value = parsed
                    cell.format_string = fmt

                if row_idx == 0:
                    cell.is_header = True
                    cell.font_bold = True

                cells.append(cell)
            table.cells.append(cells)

        if table.cells:
            table.headers = [str(c.value) for c in table.cells[0]]
            table.num_rows = len(table.cells)
            table.num_cols = max(len(r) for r in table.cells) if table.cells else 0

        return table


# ── ZIP Parser ──────────────────────────────────────────────────────────────

class ZIPParser:
    """Parses ZIP archives, extracting and processing supported files."""

    def __init__(self, process_file_callback=None):
        self.process_file_callback = process_file_callback

    def parse(self, file_path: str | Path) -> ProcessingResult:
        result = ProcessingResult()
        result.filename = Path(file_path).name
        result.file_category = "archive"
        start_time = time.time()

        try:
            result.status = ProcessingStatus.OCR_IN_PROGRESS

            from config import Config
            allowed = Config.ALLOWED_EXTENSIONS - {"zip"}

            with zipfile.ZipFile(str(file_path), 'r') as zf:
                # Create temp directory for extraction
                temp_dir = tempfile.mkdtemp(prefix="docexcel_zip_")

                try:
                    extracted_files = []
                    for info in zf.infolist():
                        if info.is_dir():
                            continue
                        ext = get_file_extension(info.filename)
                        if ext in allowed:
                            zf.extract(info, temp_dir)
                            extracted_files.append(
                                os.path.join(temp_dir, info.filename)
                            )

                    result.page_count = len(extracted_files)

                    # Process each file
                    if self.process_file_callback:
                        for extracted_path in extracted_files:
                            try:
                                sub_result = self.process_file_callback(extracted_path)
                                if sub_result and sub_result.tables:
                                    result.tables.extend(sub_result.tables)
                            except Exception as e:
                                result.warnings.append(
                                    f"Failed to process {Path(extracted_path).name}: {str(e)}"
                                )
                    else:
                        result.warnings.append(
                            f"ZIP contains {len(extracted_files)} files but no processor available"
                        )

                finally:
                    cleanup_directory(temp_dir)

            result.ocr_engine_used = "zip_parser"
            result.status = ProcessingStatus.READY_FOR_REVIEW
            result.processing_time = time.time() - start_time

        except zipfile.BadZipFile:
            result.mark_failed("Invalid or corrupted ZIP file")
        except Exception as e:
            result.mark_failed(f"ZIP parsing failed: {str(e)}")

        return result
